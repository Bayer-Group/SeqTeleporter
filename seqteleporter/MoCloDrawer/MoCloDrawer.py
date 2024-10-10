import os
import pandas as pd
import win32com.client
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def draw_a_moclo_design_with_overlap_number(presentation: Presentation,
                                            design_name: str,
                                            design_df: pd.DataFrame) -> None:

    blank_slide_layout = presentation.slide_layouts[6]
    presentation.slide_height = Inches(2.5)
    slide = presentation.slides.add_slide(blank_slide_layout)

    num_of_slots = design_df.shape[0]

    # Widths
    total_width = 8
    # Total length for drawing design is 8 inch, each slot gets a width of 8/(number of slots) inch.
    slot_width = total_width / num_of_slots

    # Heights
    function_label_height = 0.3  # set according to hand-made designs from example slides
    fusion_site_height = 0.2  # set according to hand-made designs from example slides

    # Tops
    fusion_site_top = 1
    function_label_top = fusion_site_top - (function_label_height - fusion_site_height) / 2

    # Draw a transparent box that encloses all slots, and label the design name at top-left corner
    slots_left_most = 1
    left, top, width, height = (Inches(slots_left_most - 0.5),
                                Inches(function_label_height - 0.1),
                                Inches(total_width + 1),
                                Inches(function_label_height + 1))
    design_txbox = slide.shapes.add_textbox(left, top, width, height)
    # Text
    tf = design_txbox.text_frame
    tf.text = design_name
    tf.margin_left = tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    # Line
    design_txbox.line.color.rgb = RGBColor(0, 0, 0)

    for i in range(1, num_of_slots + 1):
        # print(design_df)
        fs_5 = design_df.loc[i, "5' Fusion Site"]
        overlap5_len = int(design_df.loc[i, "5' Fusion Site - Base Overlap No."])

        if overlap5_len == 0:
            design_df.loc[i, 'fusion_site_5_no_overlap'] = fs_5
        else:
            design_df.loc[i, 'fusion_site_5_no_overlap'] = fs_5[:-overlap5_len]
        design_df.loc[i, 'fusion_site_5_overlap'] = fs_5[:overlap5_len]

        fs_3 = design_df.loc[i, "3' Fusion Site"]
        overlap3_len = int(design_df.loc[i, "3' Fusion Site - Base Overlap No."])
        if overlap3_len == 0:
            design_df.loc[i, 'fusion_site_3_no_overlap'] = fs_3
        else:
            design_df.loc[i, 'fusion_site_3_no_overlap'] = fs_3[:-overlap3_len]
        design_df.loc[i, 'fusion_site_3_overlap'] = fs_3[:overlap3_len]

        # dynamic Widths
        fusion_site_5_overlap_width = 0.1 * overlap5_len
        fusion_site_5_no_overlap_width = 0.1 * (len(fs_5) - overlap5_len)
        fusion_site_3_overlap_width = 0.1 * overlap3_len
        fusion_site_3_no_overlap_width = 0.1 * (len(fs_3) - overlap3_len)
        function_label_width = slot_width - fusion_site_5_no_overlap_width - fusion_site_3_no_overlap_width - 0.5

        # Lefts
        fusion_site_5_no_overlap_left = slots_left_most + slot_width * (i - 1)
        fusion_site_5_overlap_left = fusion_site_5_no_overlap_left + fusion_site_5_no_overlap_width
        function_label_left = fusion_site_5_overlap_left
        fusion_site_3_overlap_left = function_label_left + function_label_width - fusion_site_3_overlap_width
        fusion_site_3_no_overlap_left = function_label_left + function_label_width

        # Draw Textbox Shapes
        fusion_site_5_no_overlap_txbox = slide.shapes.add_textbox(
            Inches(fusion_site_5_no_overlap_left), Inches(fusion_site_top), Inches(fusion_site_5_no_overlap_width),
            Inches(fusion_site_height)
        )
        function_label_txbox = slide.shapes.add_textbox(
            Inches(function_label_left), Inches(function_label_top), Inches(function_label_width),
            Inches(function_label_height)
        )
        fusion_site_3_no_overlap_txbox = slide.shapes.add_textbox(
            Inches(fusion_site_3_no_overlap_left), Inches(fusion_site_top), Inches(fusion_site_3_no_overlap_width),
            Inches(fusion_site_height)
        )
        fusion_site_5_overlap_txbox = slide.shapes.add_textbox(
            Inches(fusion_site_5_overlap_left), Inches(fusion_site_top), Inches(fusion_site_5_overlap_width),
            Inches(fusion_site_height)
        )
        fusion_site_3_overlap_txbox = slide.shapes.add_textbox(
            Inches(fusion_site_3_overlap_left), Inches(fusion_site_top), Inches(fusion_site_3_overlap_width),
            Inches(fusion_site_height)
        )

        # Fill-in texts and colors
        slot_elements = [
            {'field': 'fusion_site_5_no_overlap', 'shape': fusion_site_5_no_overlap_txbox,
             'fill': ('rgb', RGBColor(255, 255, 0)), 'line': 'None', 'font_size': 8},
            {'field': 'fusion_site_5_overlap', 'shape': fusion_site_5_overlap_txbox,
             'fill': ('rgb', RGBColor(38, 234, 234)), 'line': 'None', 'font_size': 8},
            {'field': 'fusion_site_3_no_overlap', 'shape': fusion_site_3_no_overlap_txbox,
             'fill': ('rgb', RGBColor(255, 255, 0)), 'line': 'None', 'font_size': 8},
            {'field': 'fusion_site_3_overlap', 'shape': fusion_site_3_overlap_txbox,
             'fill': ('rgb', RGBColor(38, 234, 234)), 'line': 'None', 'font_size': 8},
            {'field': 'Function(s)', 'shape': function_label_txbox, 'fill': ('rgb', RGBColor(173, 223, 245)),
             'line': ('rgb', RGBColor(0, 0, 0)), 'font_size': 12}
        ]

        for slot_element in slot_elements:
            # Add Texts
            tf = slot_element['shape'].text_frame
            tf.text = design_df.loc[i][slot_element['field']]
            # tf.margin_left = tf.margin_top = tf.margin_bottom = Inches(0.02)
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.font.size = Pt(slot_element['font_size'])
            p.alignment = PP_ALIGN.CENTER
            # Add fill color
            fill = slot_element['shape'].fill
            fill.solid()
            if slot_element['fill'][0] == 'rgb':
                fill.fore_color.rgb = slot_element['fill'][1]
            # Add line color
            line = slot_element['shape'].line
            if slot_element['line'][0] == 'rgb':
                line.color.rgb = slot_element['line'][1]


def draw_moclo_designs_with_overlap_number(input_xlsx: str, output_pptx_path: str, output_image_dir: str) -> None:
    design_df = pd.read_excel(input_xlsx, index_col='Position')
    prs = Presentation()
    design_names = []
    for idx_, df in design_df.groupby('Chain Design'):
        cleaned_df = df.loc[:, df.columns != 'Protein Complex Design'].drop_duplicates()
        design_name = '_'.join(list(cleaned_df['Function(s)']))
        draw_a_moclo_design_with_overlap_number(presentation=prs, design_name=design_name, design_df=cleaned_df)
        design_names.append(design_name)

    if os.path.isfile(output_pptx_path):  # file exists
        os.remove(output_pptx_path)
    prs.save(output_pptx_path)

    if output_image_dir:
        application = win32com.client.Dispatch("PowerPoint.Application")
        presentation_ = application.Presentations.Open(output_pptx_path)
        for i in range(0, len(design_names)):
            img_path = os.path.join(output_image_dir, design_names[i] + '.jpg')
            presentation_.Slides[i].Export(img_path, "JPG")
        application.Quit()
        # Presentation_ = None
        # Application = None
